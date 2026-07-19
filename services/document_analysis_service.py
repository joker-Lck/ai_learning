"""
文档资料分析服务 - 解析学生上传的学习资料并进行AI分析
支持格式: txt, md, pdf, docx, pptx, jpg, png
"""

import io
from datetime import datetime

from core.json_utils import safe_parse_json
from core.logger import error, info


class DocumentAnalysisService:
    """文档资料分析服务"""

    SUPPORTED_EXTENSIONS = {
        'txt', 'md', 'pdf', 'doc', 'docx', 'ppt', 'pptx',
        'xls', 'xlsx', 'csv',
        'jpg', 'jpeg', 'png', 'bmp', 'webp'
    }

    MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
    MAX_TOTAL_SIZE = 30 * 1024 * 1024  # 30MB
    MAX_FILES = 10

    def __init__(self):
        info("文档资料分析服务初始化完成")

    def analyze_documents(self, files: list[dict], user_context: dict | None = None) -> dict:
        """
        分析多个文档资料

        Args:
            files: [{"filename": str, "content": bytes, "size": int}]
            user_context: {"subject": str, "topic": str, "difficulty": str}

        Returns:
            分析结果
        """
        info(f"开始分析 {len(files)} 个文档")

        # 1. 解析所有文件为文本
        parsed_files = []
        total_size = 0
        for f in files:
            if f["size"] > self.MAX_FILE_SIZE:
                parsed_files.append({
                    "filename": f["filename"],
                    "status": "error",
                    "error": "文件大小超过限制(最大10MB)",
                    "text": ""
                })
                continue

            total_size += f["size"]
            if total_size > self.MAX_TOTAL_SIZE:
                parsed_files.append({
                    "filename": f["filename"],
                    "status": "error",
                    "error": "总文件大小超过限制(最大30MB)",
                    "text": ""
                })
                continue

            text = self._parse_file(f["filename"], f["content"])
            parsed_files.append({
                "filename": f["filename"],
                "status": "success" if text and not text.startswith("[") else "warning",
                "text": text or "",
                "char_count": len(text) if text else 0
            })

        successful = [f for f in parsed_files if f["status"] in ("success", "warning")]
        if not successful:
            return {
                "success": False,
                "message": "所有文件解析失败",
                "data": {"files": parsed_files}
            }

        # 2. 合并文本内容(限制总长度)
        combined_text = self._combine_texts(successful, max_chars=12000)

        # 3. AI分析
        subject = (user_context or {}).get("subject", "")
        topic = (user_context or {}).get("topic", "")
        difficulty = (user_context or {}).get("difficulty", "")

        analysis = self._ai_analyze(combined_text, parsed_files, subject, topic, difficulty)

        # 4. 构建结果
        result = {
            "success": True,
            "message": f"成功分析 {len(successful)}/{len(files)} 个文件",
            "data": {
                "files": [{
                    "filename": f["filename"],
                    "status": f["status"],
                    "char_count": f.get("char_count", 0),
                    "error": f.get("error", "")
                } for f in parsed_files],
                "analysis": analysis,
                "analyzed_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            }
        }

        info(f"文档分析完成, 成功: {len(successful)}/{len(files)}")
        return result

    def _parse_file(self, filename: str, content: bytes) -> str | None:
        """解析单个文件为文本"""
        ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

        try:
            if ext in ('txt', 'md'):
                return content.decode('utf-8', errors='ignore')

            elif ext in ('doc', 'docx'):
                return self._parse_docx(content)

            elif ext == 'pdf':
                return self._parse_pdf(content)

            elif ext in ('ppt', 'pptx'):
                return self._parse_pptx(content)

            elif ext in ('xls', 'xlsx', 'csv'):
                return self._parse_excel(content, ext)

            elif ext in ('jpg', 'jpeg', 'png', 'bmp', 'webp'):
                return f"[图片文件: {filename} - 需要视觉模型识别]"

            else:
                return f"[不支持的文件格式: .{ext}]"

        except Exception as e:
            error(f"解析文件 {filename} 失败: {e!s}")
            return f"[解析失败: {e!s}]"

    def _parse_docx(self, content: bytes) -> str:
        """解析 Word 文档"""
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs)
        except ImportError:
            return "[需要安装 python-docx: pip install python-docx]"
        except Exception as e:
            return f"[Word解析失败: {e!s}]"

    def _parse_pdf(self, content: bytes) -> str:
        """解析 PDF 文档（支持文字版和扫描版）"""
        text = ""
        page_count = 0

        # 优先用 PyMuPDF（支持更多 PDF 格式）
        try:
            import fitz
            doc = fitz.open(stream=content, filetype="pdf")
            page_count = len(doc)
            for page in doc:
                page_text = page.get_text() or ""
                cleaned = ''.join(c for c in page_text if c.isprintable() or c in '\n\r\t')
                if cleaned.strip():
                    text += cleaned + "\n"
            doc.close()
            info(f"[PDF] PyMuPDF: {len(text)} chars from {page_count} pages")
        except ImportError:
            info("[PDF] PyMuPDF not available")
        except Exception as e:
            info(f"[PDF] PyMuPDF error: {e}")

        # 如果 PyMuPDF 没提取到文本，尝试 PyPDF2
        if not text.strip():
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(io.BytesIO(content))
                page_count = len(reader.pages)
                for page in reader.pages:
                    page_text = page.extract_text() or ""
                    cleaned = ''.join(c for c in page_text if c.isprintable() or c in '\n\r\t')
                    if cleaned.strip():
                        text += cleaned + "\n"
                info(f"[PDF] PyPDF2: {len(text)} chars from {page_count} pages")
            except Exception as e:
                info(f"[PDF] PyPDF2 error: {e}")

        # 判断是否需要 OCR：文本量太少（平均每页<100字符）视为扫描版
        text_len = len(text.strip())
        need_ocr = False
        if text_len == 0:
            need_ocr = True
            info("[PDF] 无文本，需要 OCR")
        elif page_count > 0 and text_len / page_count < 100:
            need_ocr = True
            info(f"[PDF] 文本偏少（{text_len}/{page_count}页），需要 OCR")

        if need_ocr:
            info("[PDF] 开始 AI OCR...")
            ocr_text = self._ocr_pdf_with_ai(content)
            info(f"[PDF] OCR 完成: {len(ocr_text)} chars")
            if len(ocr_text) > text_len:
                text = ocr_text
                info("[PDF] 使用 OCR 结果")

        info(f"[PDF] 最终: {len(text)} chars")
        return text.strip()

    def _ocr_pdf_with_ai(self, content: bytes, max_pages: int = 50) -> str:
        """用 AI 视觉模型识别扫描版 PDF（逐页渲染 + AI 识别）"""
        try:
            import fitz
            import base64
            from services.spark_client import spark_client
        except ImportError:
            error("[OCR] PyMuPDF 未安装")
            return "[扫描版 PDF 需要安装 PyMuPDF: pip install PyMuPDF]"

        try:
            doc = fitz.open(stream=content, filetype="pdf")
            total_pages = len(doc)
            info(f"[OCR] 总页数: {total_pages}, 最大处理: {max_pages}")

            all_text = []
            # 采样策略：超过 max_pages 时均匀采样
            if total_pages > max_pages:
                step = total_pages // max_pages
                page_indices = list(range(0, total_pages, step))[:max_pages]
                info(f"[OCR] 采样: {len(page_indices)} 页 (每 {step} 页)")
            else:
                page_indices = list(range(total_pages))
                info(f"[OCR] 处理全部 {len(page_indices)} 页")

            success_count = 0
            fail_count = 0
            for i, idx in enumerate(page_indices):
                try:
                    page = doc[idx]
                    pix = page.get_pixmap(dpi=150)
                    img_bytes = pix.tobytes("png")
                    img_b64 = "data:image/png;base64," + base64.b64encode(img_bytes).decode()

                    prompt = "请提取这张图片中的所有文字内容。只输出文字。"
                    result = spark_client.chat_with_image(
                        prompt=prompt,
                        image_b64=img_b64,
                        max_tokens=4000
                    )
                    if result and len(result) > 10 and not result.startswith("["):
                        all_text.append(result)
                        success_count += 1
                    else:
                        fail_count += 1
                        if i < 3:
                            info(f"[OCR] 第{idx+1}页: 失败 result={repr(result[:80]) if result else 'None'}")

                    if (i + 1) % 10 == 0:
                        info(f"[OCR] 进度: {i+1}/{len(page_indices)}, 成功: {success_count}")

                except Exception as e:
                    fail_count += 1
                    if i < 3:
                        error(f"[OCR] 第{idx+1}页异常: {e}")
                    continue

            doc.close()
            result_text = "\n\n".join(all_text)
            info(f"[OCR] 完成: 成功 {success_count}, 失败 {fail_count}, 文本 {len(result_text)} chars")
            return result_text if result_text else "[扫描版 PDF AI 识别未提取到文本]"

        except Exception as e:
            error(f"[OCR] 异常: {e}")
            return f"[扫描版 PDF 识别失败: {e}]"

    def _parse_pptx(self, content: bytes) -> str:
        """解析 PPT 文档"""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slides_text.append(shape.text.strip())
            return "\n\n".join(slides_text)
        except ImportError:
            return "[需要安装 python-pptx: pip install python-pptx]"
        except Exception as e:
            return f"[PPT解析失败: {e!s}]"

    def _parse_excel(self, content: bytes, ext: str) -> str:
        """解析 Excel/CSV 文件"""
        try:
            import pandas as pd
            if ext == 'csv':
                df = pd.read_csv(io.BytesIO(content))
            elif ext == 'xls':
                df = pd.read_excel(io.BytesIO(content), engine='xlrd')
            else:
                df = pd.read_excel(io.BytesIO(content))
            # 转为文本表格格式
            lines = []
            for _, row in df.iterrows():
                row_text = ' | '.join(str(v) for v in row.values if pd.notna(v))
                if row_text.strip():
                    lines.append(row_text)
            return "\n".join(lines)
        except ImportError:
            return "[需要安装 pandas: pip install pandas openpyxl xlrd]"
        except Exception as e:
            return f"[Excel解析失败: {e!s}]"

    def _combine_texts(self, files: list[dict], max_chars: int = 12000) -> str:
        """合并多个文件的文本内容"""
        parts = []
        total = 0

        for f in files:
            text = f.get("text", "")
            if not text or text.startswith("["):
                continue

            remaining = max_chars - total
            if remaining <= 0:
                break

            if len(text) > remaining:
                text = text[:remaining] + "\n...(内容截断)"

            parts.append(f"===== 文件: {f['filename']} =====\n{text}")
            total += len(text)

        return "\n\n".join(parts)

    def _ai_analyze(self, text: str, files_info: list[dict],
                    subject: str, topic: str, difficulty: str) -> dict:
        """调用AI分析文档内容"""

        file_list = ", ".join([f["filename"] for f in files_info if f["status"] != "error"])

        prompt = f"""你是一位专业的教育分析师。请分析以下学生上传的学习资料，生成详细的学习效果分析报告。

文件列表: {file_list}
{f'学科: {subject}' if subject else ''}
{f'主题: {topic}' if topic else ''}
{f'难度: {difficulty}' if difficulty else ''}

文档内容:
{text[:10000]}

请从以下维度进行分析，并输出JSON格式(只输出JSON):

{{
    "knowledge_overview": {{
        "total_knowledge_points": 数字,
        "main_topics": ["主题1", "主题2", "主题3"],
        "coverage_summary": "资料覆盖范围概述(2-3句话)"
    }},
    "knowledge_points": [
        {{
            "name": "知识点名称",
            "importance": "核心/重要/辅助",
            "mastery_hint": "已掌握/需巩固/薄弱",
            "description": "简要说明"
        }}
    ],
    "strengths": ["学生已掌握的优势领域1", "优势2", "优势3"],
    "weaknesses": ["需要加强的薄弱环节1", "薄弱环节2", "薄弱环节3"],
    "learning_gaps": [
        {{
            "gap": "知识缺口描述",
            "related_topics": ["相关主题"],
            "suggestion": "补充建议"
        }}
    ],
    "difficulty_assessment": {{
        "overall_level": "入门/中级/高级",
        "reasoning": "判断依据"
    }},
    "study_recommendations": [
        {{
            "priority": "高/中/低",
            "action": "具体学习行动建议",
            "resources": "推荐资源类型(document/quiz/video/code_case)"
        }}
    ],
    "overall_score": 75,
    "summary": "用Markdown格式撰写的综合分析报告(200-300字)，包含:资料质量评估、知识掌握情况、个性化学习建议"
}}
"""

        try:
            from services.spark_client import spark_client
            response = spark_client.simple(prompt, max_tokens=3000)
            analysis = safe_parse_json(response)

            if not analysis:
                return self._fallback_analysis(files_info)

            return analysis

        except Exception as e:
            error(f"AI文档分析失败: {e!s}")
            return self._fallback_analysis(files_info)

    def _fallback_analysis(self, files_info: list[dict]) -> dict:
        """降级分析方案"""
        successful = [f for f in files_info if f["status"] != "error"]
        total_chars = sum(f.get("char_count", 0) for f in successful)

        return {
            "knowledge_overview": {
                "total_knowledge_points": 0,
                "main_topics": [],
                "coverage_summary": f"共解析 {len(successful)} 个文件，约 {total_chars} 字符。AI分析暂时不可用，请稍后重试。"
            },
            "knowledge_points": [],
            "strengths": [],
            "weaknesses": [],
            "learning_gaps": [],
            "difficulty_assessment": {
                "overall_level": "未知",
                "reasoning": "AI分析不可用"
            },
            "study_recommendations": [],
            "overall_score": 0,
            "summary": "AI分析服务暂时不可用，请稍后重试。"
        }


# 全局实例
document_analysis_service = DocumentAnalysisService()
